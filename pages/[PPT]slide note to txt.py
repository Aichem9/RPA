import streamlit as st
import pptx

# 페이지 설명 부분
st.title("🖨 PPT 슬라이드 노트 reader")


st.info('###### 언제 사용하나요?\nPPT 슬라이드별 노트를 한꺼번에 출력해서 발표를 준비하고 싶을 때')
st.warning('###### 어떻게 해결하나요?\nPPT 파일 업로드 ➡ 슬라이드 노트 한꺼번에 출력')


st.write("※ 업로드하신 PPT는 따로 저장되지 않으므로 걱정하지 않으셔도 됩니다. ")
def get_slide_notes(pptx_file):
    prs = pptx.Presentation(pptx_file)
    slide_notes = []

    for slide in prs.slides:
        notes_slide = slide.notes_slide
        if notes_slide is not None:
            notes_text = notes_slide.notes_text_frame.text
            slide_notes.append(notes_text)
    return slide_notes

# Streamlit 앱 설정
#st.set_page_config(page_title="PPTX Viewer")

# 업로드된 파일 가져오기
st.write("### pptx 파일 업로드")
uploaded_file = st.file_uploader("파일을 먼저 업로드하세요.", type=["pptx"])
separator = st.text_input("슬라이드별로 구분할 단어를 입력해주세요", value= '슬라이드')

if uploaded_file is not None:
    slide_notes = get_slide_notes(uploaded_file)

    # "파싱 시작" 버튼 추가
    if st.checkbox("슬라이드 노트 출력하기"):
        for i, note in enumerate(slide_notes, start=1):
            st.markdown("#### {} {}".format(separator, i))
            st.write(note)

    elif st.checkbox("구분자 없이 출력할게요"):
        for i, note in enumerate(slide_notes, start=1):
            st.write(note)
